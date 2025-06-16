import os
import json
import re
import requests
import numpy as np
import fitz
import pandas as pd
from pptx import Presentation
from io import BytesIO
from docx import Document  # python-docx

# external helpers
from query_handler import (
    embed_query_sentence,  # returns a 768-d numpy vector
    embedding_model,       # SentenceTransformer already loaded
    query_openai,
    search_topk,
)

#  MIME → export map (Google-native ↦ local format)
EXPORT_MIME = {
    "google_doc": "text/plain",           # Docs → TXT
    "spreadsheet": "text/csv",            # Sheets → CSV
    "presentation": "application/pdf",    # Slides → PDF
}

#  Logical type sets
TEXT_TYPES  = {
    "pdf", "google_doc", "spreadsheet", "presentation", "pptx",
    "text", "xlsx", "docx", "csv", "google_slide", "google_sheet",
}
MEDIA_TYPES = {"image", "video", "audio"}

#  Utilities
def is_text_type(ftype: str) -> bool:
    return ftype in TEXT_TYPES

def icon_for(ftype: str) -> str:
    return {
        "pdf": "📄", "google_doc": "📄", "text": "📄", "docx": "📄",
        "spreadsheet": "📊", "presentation": "📽️",
        "image": "🖼️", "video": "🎞️", "audio": "🎧",
        "folder": "📁",
    }.get(ftype, "📦")

#  Simple chunk-ranker (semantic similarity)
def rank_chunks(query: str, chunks: list[str], top_k: int = 5) -> list[str]:
    """
    Embed *query* and *chunks*, return the top-k most similar chunks.
    """
    if not chunks:
        return []
    q_vec  = embed_query_sentence(query)                     # (768,)
    c_vecs = embedding_model.encode(chunks, convert_to_numpy=True)
    sims   = np.dot(c_vecs, q_vec) / (
        (np.linalg.norm(c_vecs, axis=1) + 1e-8) * np.linalg.norm(q_vec)
    )
    idx = sims.argsort()[::-1][:top_k]
    return [chunks[i] for i in idx]

# Download + export
def download_file(file_id: str, user_id: str, token: str) -> str | None:
    url  = f"https://www.googleapis.com/drive/v3/files/{file_id}?alt=media"
    hdr  = {"Authorization": f"Bearer {token}"}
    r    = requests.get(url, headers=hdr, stream=True)
    if r.status_code != 200:
        print("❌ download failed:", r.text)
        return None
    path = f"user_data/{user_id}/downloads/{file_id}"
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "wb") as f:
        for chunk in r.iter_content(8192):
            f.write(chunk)
    return path

def export_google_file(file_id: str, logical_type: str, user_id: str, token: str) -> str | None:
    mime = EXPORT_MIME.get(logical_type)
    if not mime:
        return None
    url  = f"https://www.googleapis.com/drive/v3/files/{file_id}/export"
    hdr  = {"Authorization": f"Bearer {token}"}
    r    = requests.get(url, headers=hdr, params={"mimeType": mime})
    if r.status_code != 200:
        print("❌ export failed:", r.text)
        return None
    ext_map = {"text/plain": "txt", "text/csv": "csv"}
    ext = ext_map.get(mime, "pdf" if "pdf" in mime else "xlsx")
    path = f"user_data/{user_id}/downloads/{file_id}.{ext}"
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "wb") as f:
        f.write(r.content)
    return path

# Extractors for each type
def extract_text_from_pdf(path: str) -> str:
    doc = fitz.open(path)
    return "\n".join(p.get_text() for p in doc)

def extract_text_from_docx(path: str) -> str:
    try:
        doc = Document(path)
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.extend(cell.text for cell in row.cells if cell.text)
        return "\n".join(t for t in parts if t.strip())
    except Exception as e:
        return f"(⚠️ DOCX read error: {e})"

def extract_text_from_csv(path: str) -> str:
    try:
        df = pd.read_csv(path, nrows=20)
        return df.to_string(index=False)
    except Exception as e:
        return f"(⚠️ CSV read error: {e})"

def extract_text_from_excel(path: str) -> str:
    try:
        df = pd.read_excel(path, engine="openpyxl", nrows=20)
        return df.to_string(index=False)
    except Exception as e:
        return f"(⚠️ Excel read error: {e})"

def extract_text_from_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
        return "\n".join(
            shape.text for slide in prs.slides
            for shape in slide.shapes if hasattr(shape, "text")
        )
    except Exception as e:
        return f"(⚠️ PPTX read error: {e})"

# Type-aware processing
def process_file(path: str, logical_type: str) -> str:
    if logical_type == "pdf":
        return extract_text_from_pdf(path)
    if logical_type in {"text", "google_doc"}:
        return open(path, "r", encoding="utf-8").read()
    if logical_type == "docx":
        return extract_text_from_docx(path)
    if logical_type in {"spreadsheet", "xlsx"}:
        return extract_text_from_excel(path)
    if logical_type == "csv":
        return extract_text_from_csv(path)
    if logical_type in {"pptx", "presentation"}:
        return extract_text_from_pptx(path)
    return "⚠️ Unsupported file type"

# Chunkers
def chunk_text(txt: str, size: int = 500, overlap: int = 100):
    words, out, i = txt.split(), [], 0
    while i < len(words):
        out.append(" ".join(words[i : i + size]))
        i += size - overlap
    return out

# Download and extract logic
def _handle_google_doc(doc, uid, token):
    exported = export_google_file(doc["id"], "google_doc", uid, token)
    if exported:
        return exported, "text"
    return _handle_uploaded_docx(doc, uid, token)

def _handle_uploaded_docx(doc, uid, token):
    p = download_file(doc["id"], uid, token)
    return (p, "docx") if p else (None, None)

def _handle_google_sheet(doc, uid, token):
    mime = doc["raw"]["mimeType"]
    if mime.startswith("application/vnd.google-apps."):
        csv = export_google_file(doc["id"], "spreadsheet", uid, token)
        return (csv, "csv") if csv else (None, None)
    xlsx = download_file(doc["id"], uid, token)
    return (xlsx, "xlsx") if xlsx else (None, None)

def download_and_extract_top_files(docs, uid, token):
    out = []
    for d in docs:
        path, ltype = None, None
        if d["type"] == "google_doc":
            path, ltype = _handle_google_doc(d, uid, token)
        elif d["type"] in {"google_sheet", "spreadsheet", "xlsx"}:
            path, ltype = _handle_google_sheet(d, uid, token)
        elif d["type"] == "docx":
            path, ltype = _handle_uploaded_docx(d, uid, token)
        else:
            if is_text_type(d["type"]):
                path = download_file(d["id"], uid, token)
                ltype = d["type"]

        if not path:
            continue
        text = process_file(path, ltype)
        if text.strip():
            out.append({"doc": d, "chunks": chunk_text(text)})
    return out

# Rag helper
def generate_response_with_context(query, context_chunks, history=None):
    hist = ""
    if history:
        hist = "\n\n### Conversation so far ###\n" + "\n\n".join(
            f"USER: {h['q']}\nASSISTANT: {h['a']}" for h in history[-5:]
        )
    context = "\n---\n".join(context_chunks)
    prompt = (
        "You are a helpful assistant. Answer the query using ONLY the context below."
        f"{hist}\n\n### Context ###\n{context}\n\n### Query ###\n{query}\n\n### Answer ###"
    )
    return query_openai(prompt, max_tokens=300)

# Folder helper
def list_folder_children(fid, token, limit=10):
    q   = f"'{fid}' in parents and trashed=false"
    hdr = {"Authorization": f"Bearer {token}"}
    r   = requests.get(
        "https://www.googleapis.com/drive/v3/files",
        headers=hdr,
        params={"q": q, "pageSize": limit, "fields": "files(id,name,mimeType,webViewLink)"},
    )
    if r.status_code != 200:
        return []
    return [
        {"name": f["name"], "type": f["mimeType"], "link": f.get("webViewLink")}
        for f in r.json().get("files", [])
    ]

# Final response logic
def generate_final_response(
    user_query: str,
    user_id: str,
    results: list[dict],
    access_token: str,
    history=None,
):
    if not results:
        return {
            "answer": "I couldn't find anything relevant about that query in your Google Drive.",
            "sources": [],
        }

    first = results[0]

    def enrich(doc):
        raw = doc.get("raw", {})
        if thumb := raw.get("thumbnailLink"):
            x = dict(doc)
            x["thumb"] = thumb
            return x
        return doc

    # Folder
    if first["type"] == "folder":
        kids    = list_folder_children(first["id"], access_token, 15)
        listing = "\n".join(f"{icon_for(c['type'])} {c['name']}" for c in kids) or "*(folder is empty)*"
        return {"answer": f"📁 Folder **{first['name']}** contents:\n\n{listing}", "sources": [enrich(first)]}

    # Media
    if first["type"] in MEDIA_TYPES:
        return {"answer": f"I found a {first['type']} named **{first['name']}**. Need anything else?",
                "sources": [enrich(first)]}

    # Text Branch
    text_docs  = [d for d in results if is_text_type(d["type"])]
    other_docs = [d for d in results if not is_text_type(d["type"])]

    extracted = download_and_extract_top_files(text_docs, user_id, access_token)

    context_parts = []
    for e in extracted:
        best = rank_chunks(user_query, e["chunks"], top_k=5)
        if best:
            header = f"### {e['doc']['name']}"
            context_parts.append(header + "\n" + "\n".join(best))

    if context_parts:
        context_chunks = "\n\n".join(context_parts)
        answer = generate_response_with_context(user_query, [context_chunks], history)
        if other_docs:
            extra = ", ".join(d["name"] for d in other_docs[:3])
            answer += f"\n\n(Also matched media files: {extra})"
    else:
        names  = "\n".join(f"- {d['name']}" for d in results[:3])
        answer = f"I found these files:\n{names}\n\nLet me know which one to explore."

    sources = [enrich(e["doc"]) for e in extracted] + [enrich(d) for d in other_docs]
    return {"answer": answer, "sources": sources}

# -------------------------------------------------------TESTING--------------------------------------
# if __name__ == "__main__":
#     user_id = ""
#     token_path = f"user_data/{user_id}/tokens.json"
#     if not os.path.exists(token_path):
#         print("❌ Token file missing. Run the login flow first.")
#         exit()

#     access_token = json.load(open(token_path))["access_token"]

#     test_queries = ["Please summarize Gatech Plan Soham Agarwal for me"]

#     for q in test_queries:
#         print("\n" + "=" * 30)
#         print("🧠 Query:", q)
#         hits = search_topk(user_id, q, top_k=5)
#         print("🔎 Search hits:", len(hits))
#         resp = generate_final_response(q, user_id, hits, access_token)
#         print("🗣️", resp["answer"])
#         for s in resp["sources"]:
#             print("-", s["name"])
